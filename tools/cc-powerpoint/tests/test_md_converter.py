"""Tests for cc-powerpoint markdown converter (PPTX -> Markdown)."""

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.md_converter import convert_pptx_to_markdown


def _create_pptx(path: Path, slides_content: list[list[str]]) -> Path:
    """Helper: create a test PPTX with given slide content.

    Each slide gets a title (first item) and body bullets (remaining items).
    """
    prs = Presentation()
    for content in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(slide_layout)

        if content:
            slide.shapes.title.text = content[0]
            if len(content) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.text = content[1]
                for item in content[2:]:
                    p = tf.add_paragraph()
                    p.text = item

    prs.save(str(path))
    return path


class TestConvertPptxToMarkdown:
    """Tests for convert_pptx_to_markdown."""

    def test_basic_slide(self, tmp_path):
        pptx_path = _create_pptx(
            tmp_path / "test.pptx",
            [["My Title", "Bullet one", "Bullet two"]],
        )
        output = tmp_path / "test.md"

        md = convert_pptx_to_markdown(pptx_path, output)

        assert "My Title" in md
        assert "Bullet one" in md
        assert "Bullet two" in md

    def test_multiple_slides_separated_by_hr(self, tmp_path):
        pptx_path = _create_pptx(
            tmp_path / "test.pptx",
            [["Slide 1"], ["Slide 2"], ["Slide 3"]],
        )
        output = tmp_path / "test.md"

        md = convert_pptx_to_markdown(pptx_path, output)

        assert md.count("---") == 2  # 3 slides = 2 separators

    def test_title_is_heading(self, tmp_path):
        pptx_path = _create_pptx(
            tmp_path / "test.pptx",
            [["Important Title"]],
        )
        output = tmp_path / "test.md"

        md = convert_pptx_to_markdown(pptx_path, output)

        assert "# Important Title" in md

    def test_no_images_no_directory(self, tmp_path):
        pptx_path = _create_pptx(
            tmp_path / "test.pptx",
            [["Title Only"]],
        )
        output = tmp_path / "test.md"

        md = convert_pptx_to_markdown(pptx_path, output)

        images_dir = tmp_path / "test_images"
        assert not images_dir.exists()

    def test_image_extraction(self, tmp_path):
        # Create PPTX with an image
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Create a tiny PNG-like blob
        img_path = tmp_path / "tiny.png"
        # Minimal valid PNG (1x1 transparent pixel)
        import struct
        import zlib

        def _make_png():
            sig = b"\x89PNG\r\n\x1a\n"
            # IHDR
            ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
            ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
            ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
            # IDAT
            raw = zlib.compress(b"\x00\xFF\x00\x00")
            idat_crc = zlib.crc32(b"IDAT" + raw) & 0xFFFFFFFF
            idat = struct.pack(">I", len(raw)) + b"IDAT" + raw + struct.pack(">I", idat_crc)
            # IEND
            iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
            iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
            return sig + ihdr + idat + iend

        img_path.write_bytes(_make_png())
        slide.shapes.add_picture(str(img_path), Inches(1), Inches(1))

        pptx_path = tmp_path / "with_img.pptx"
        prs.save(str(pptx_path))
        output = tmp_path / "with_img.md"

        md = convert_pptx_to_markdown(pptx_path, output)

        images_dir = tmp_path / "with_img_images"
        assert images_dir.exists()
        assert any(images_dir.iterdir())
        assert "with_img_images/" in md

    def test_output_ends_with_newline(self, tmp_path):
        pptx_path = _create_pptx(tmp_path / "test.pptx", [["Title"]])
        output = tmp_path / "test.md"

        md = convert_pptx_to_markdown(pptx_path, output)

        assert md.endswith("\n")

    def test_speaker_notes(self, tmp_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Noted Slide"
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "Speaker notes here"

        pptx_path = tmp_path / "notes.pptx"
        prs.save(str(pptx_path))
        output = tmp_path / "notes.md"

        md = convert_pptx_to_markdown(pptx_path, output)

        assert "> Speaker notes here" in md
