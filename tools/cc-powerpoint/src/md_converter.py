"""Convert PowerPoint (PPTX) files to Markdown with image extraction."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from cc_shared.image_extractor import ExtractedImage, save_extracted_images
except ImportError:
    import sys
    import os
    if hasattr(sys, '_MEIPASS'):
        sys.path.insert(0, os.path.join(sys._MEIPASS, 'cc_shared'))
    else:
        sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent / "cc_shared"))
    from image_extractor import ExtractedImage, save_extracted_images


def convert_pptx_to_markdown(
    input_path: Path,
    output_path: Path,
) -> str:
    """Convert a PowerPoint presentation to Markdown.

    Outputs slides separated by ``---`` (matching the format that cc-powerpoint
    accepts as input), with images extracted to a sibling directory.

    Args:
        input_path: Path to the ``.pptx`` file.
        output_path: Path to the output ``.md`` file.

    Returns:
        Markdown string.
    """
    prs = Presentation(str(input_path))
    images: list[ExtractedImage] = []
    image_placeholder_map: dict[int, int] = {}  # slide_image_counter -> images index
    slides_md: list[str] = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_lines: list[str] = []
        slide_image_counter = 0

        # Process shapes in order (top to bottom, left to right)
        shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))

        for shape in shapes:
            # Image shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                img = shape.image
                ext = img.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                if ext == "svg+xml":
                    ext = "svg"

                img_idx = len(images)
                name = f"slide{slide_idx + 1}_image_{slide_image_counter + 1:03d}.{ext}"
                images.append(ExtractedImage(
                    data=img.blob,
                    extension=ext,
                    original_name=name,
                    alt_text="",
                ))
                image_placeholder_map[img_idx] = img_idx
                slide_lines.append(f"![]({{__pptx_image_{img_idx}__}})")
                slide_image_counter += 1
                continue

            # Text shapes (titles, body text, etc.)
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame

            # Detect if this is a title shape
            is_title = False
            if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                ph_idx = shape.placeholder_format.idx
                # 0 = title, 1 = center title, 13 = title
                if ph_idx in (0, 1, 13):
                    is_title = True

            for para in text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                if is_title:
                    slide_lines.insert(0, f"# {text}")
                    is_title = False  # Only first paragraph is title
                else:
                    # Check indentation level for sub-bullets
                    level = para.level if hasattr(para, "level") else 0
                    if level > 0:
                        indent = "  " * level
                        slide_lines.append(f"{indent}- {text}")
                    elif _looks_like_bullet(para):
                        slide_lines.append(f"- {text}")
                    else:
                        slide_lines.append(text)

        # Process notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_lines.append("")
                slide_lines.append(f"> {notes_text}")

        slides_md.append("\n".join(slide_lines))

    # Join slides with --- separator
    markdown = "\n\n---\n\n".join(slides_md)

    # Save images and replace placeholders
    if images:
        path_map = save_extracted_images(images, output_path)
        for img_idx, img in enumerate(images):
            placeholder = f"{{__pptx_image_{img_idx}__}}"
            rel_path = path_map.get(img.original_name, placeholder)
            markdown = markdown.replace(placeholder, rel_path)

    return markdown.strip() + "\n"


def _looks_like_bullet(paragraph) -> bool:
    """Heuristic: does this paragraph look like a bullet point?

    Checks if the paragraph has bullet XML formatting.
    """
    pPr = paragraph._p.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}pPr"
    )
    if pPr is not None:
        buNone = pPr.find(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone"
        )
        if buNone is not None:
            return False
        # If there's any bullet definition, treat as bullet
        for child in pPr:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag.startswith("bu") and tag != "buNone":
                return True
    return False
